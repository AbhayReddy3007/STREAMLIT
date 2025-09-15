import copy
import re
import datetime
import streamlit as st
import docx
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
import google.generativeai as genai

# ---------------- CONFIG ----------------
API_KEY = "AIzaSyASAUFBojVTrN6wFN2JormPrL2sZWQZGWA"  # 🔑 Hardcode your Gemini API key
MODEL_NAME = "gemini-2.0-flash"
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel(MODEL_NAME)

st.set_page_config(page_title="AI Productivity Suite", layout="wide")
st.title("📑📄 PPT & DOC Generator + Chat with Document")

# ---------------- HELPERS ----------------
def call_gemini(prompt: str) -> str:
    resp = model.generate_content(prompt)
    return resp.text.strip()

def extract_slide_count(description: str, default: int = 5) -> int:
    m = re.search(r"(\d+)\s*(slides?|sections?|pages?)", description, re.IGNORECASE)
    if m:
        total = int(m.group(1))
        return max(1, total - 1)
    return default - 1

def generate_outline_from_desc(description: str, num_items: int, mode: str = "ppt"):
    if mode == "ppt":
        prompt = f"""Create a PowerPoint outline on: {description}.
Generate exactly {num_items} content slides (excluding title slide).
Format strictly like this:
Slide 1: <Title>
• Main point
- Sub-point
"""
    else:
        prompt = f"""Create a detailed Document outline on: {description}.
Generate exactly {num_items} sections (treat each section as roughly one page).
Each section should have:
- A section title
- 2–3 descriptive paragraphs (5–7 sentences each) of full prose, not bullets.
Format strictly like this:
Section 1: <Title>
<Paragraph 1>
<Paragraph 2>
"""
    points_text = call_gemini(prompt)
    return parse_points(points_text)

def parse_points(points_text: str):
    points = []
    current_title, current_content = None, []
    lines = [re.sub(r"[#*>`]", "", ln).strip() for ln in points_text.splitlines()]

    for line in lines:
        if not line:
            continue
        m = re.match(r"^\s*(Slide|Section)\s*(\d+)\s*:\s*(.+)$", line, re.IGNORECASE)
        if m:
            if current_title:
                points.append({"title": current_title, "description": "\n\n".join(current_content)})
            current_title, current_content = m.group(3).strip(), []
            continue
        if line.startswith("•"):
            current_content.append("• " + line.lstrip("•").strip())  # main bullet
        elif line.startswith("-"):
            current_content.append("- " + line.lstrip("-").strip())  # sub-point
        else:
            current_content.append(line.strip())

    if current_title:
        points.append({"title": current_title, "description": "\n\n".join(current_content)})
    return points

def extract_text(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        text_parts = [page.get_text("text") for page in doc]
        return "\n".join(text_parts)
    if name.endswith(".docx"):
        d = docx.Document(uploaded_file)
        return "\n".join(p.text for p in d.paragraphs)
    if name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8", errors="ignore")
    return ""

def summarize_long_text(full_text: str) -> str:
    return call_gemini(f"Summarize this text clearly:\n\n{full_text}")

def generate_title(summary: str) -> str:
    prompt = f"""
    Read the following summary and generate ONE short, clear, presentation-style title.
    - Keep it under 12 words
    - Do not give multiple options
    - Return ONLY the title, nothing else.

    Summary:
    {summary}
    """
    return call_gemini(prompt).strip()

# ---------------- FILE GENERATORS ----------------
def create_ppt(title, slides, filename="output.pptx"):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = ""

    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for slide_data in slides:
        slide = prs.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        shapes.title.text = slide_data["title"]

        text_frame = shapes.placeholders[1].text_frame
        text_frame.clear()
        for line in slide_data["description"].split("\n"):
            if not line.strip():
                continue
            if line.strip().startswith("-"):  # sub-point
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.level = 1
            else:  # main point
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.level = 0

    prs.save(filename)

def create_doc(title, sections, filename="output.docx"):
    doc = docx.Document()
    doc.add_heading(title, 0)
    for section in sections:
        doc.add_heading(section["title"], level=1)
        for para in section["description"].split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
    doc.save(filename)

# ---------------- STATE ----------------
defaults = {
    "messages": [],
    "outline": None,
    "outline_mode": None,
    "summary_text": None,
    "summary_title": None,
    "doc_chat_history": [],
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# ---------------- CHAT ----------------
for role, content in st.session_state.messages:
    with st.chat_message(role):
        st.markdown(content)

if prompt := st.chat_input("Type a message, ask for a PPT or a Document..."):
    st.session_state.messages.append(("user", prompt))
    text = prompt.lower()

    try:
        if "ppt" in text or "presentation" in text or "slides" in text:
            outline = generate_outline_from_desc(prompt, extract_slide_count(prompt), mode="ppt")
            st.session_state.outline = {"title": generate_title(prompt), "slides": outline}
            st.session_state.outline_mode = "ppt"
            st.session_state.messages.append(("assistant", "✅ Generated PPT outline! Preview below."))
        elif "doc" in text or "document" in text or "report" in text or "pages" in text:
            outline = generate_outline_from_desc(prompt, extract_slide_count(prompt), mode="doc")
            st.session_state.outline = {"title": generate_title(prompt), "sections": outline}
            st.session_state.outline_mode = "doc"
            st.session_state.messages.append(("assistant", "✅ Generated Document outline! Preview below."))
        else:
            bot_reply = call_gemini(prompt)
            st.session_state.messages.append(("assistant", bot_reply))
    except Exception as e:
        st.session_state.messages.append(("assistant", f"⚠️ Error: {e}"))
    st.rerun()

# ---------------- OUTLINE PREVIEW ----------------
if st.session_state.outline:
    mode = st.session_state.outline_mode
    outline = st.session_state.outline

    st.subheader(f"📝 Preview Outline ({mode.upper()})")
    new_title = st.text_input("📌 Edit Title", value=outline.get("title", "Untitled"), key=f"title_{mode}")

    for idx, item in enumerate(outline.get("slides", []) if mode == "ppt" else outline.get("sections", []), start=1):
        with st.expander(f"{'Slide' if mode=='ppt' else 'Section'} {idx}: {item['title']}"):
            st.markdown(item["description"])

    col1, col2 = st.columns(2)
    with col1:
        feedback = st.text_area("✏️ Feedback for outline", key=f"feedback_{mode}")
        if st.button("🔄 Apply Feedback"):
            prompt = f"Update this {mode.upper()} outline with feedback:\n\n{outline}\n\nFeedback: {feedback}"
            updated_points = parse_points(call_gemini(prompt))
            if mode == "ppt":
                st.session_state.outline = {"title": new_title, "slides": updated_points}
            else:
                st.session_state.outline = {"title": new_title, "sections": updated_points}
            st.rerun()

    with col2:
        if st.button(f"✅ Generate {mode.upper()}"):
            filename = f"{new_title.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ 'pptx' if mode=='ppt' else 'docx'}"
            if mode == "ppt":
                create_ppt(new_title, st.session_state.outline["slides"], filename)
                with open(filename, "rb") as f:
                    st.download_button("⬇️ Download PPT", f, file_name=filename,
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            else:
                create_doc(new_title, st.session_state.outline["sections"], filename)
                with open(filename, "rb") as f:
                    st.download_button("⬇️ Download DOC", f, file_name=filename,
                                       mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# ---------------- DOCUMENT UPLOAD ----------------
uploaded_file = st.file_uploader("📂 Upload a document", type=["pdf", "docx", "txt"])
if uploaded_file:
    with st.spinner("Extracting & summarizing..."):
        text = extract_text(uploaded_file)
        summary = summarize_long_text(text)
        st.session_state.summary_text = summary
        st.session_state.summary_title = generate_title(summary)  # ✅ one clean title
    st.success(f"✅ Document processed! Title: **{st.session_state.summary_title}**")

    st.markdown("💬 **Chat with your document**")
    for role, content in st.session_state.doc_chat_history:
        with st.chat_message(role):
            st.markdown(content)

    if doc_prompt := st.chat_input("Ask a question about the uploaded document..."):
        st.session_state.doc_chat_history.append(("user", doc_prompt))
        answer = call_gemini(
            f"Answer based only on this document:\n\n{st.session_state.summary_text}\n\nQ: {doc_prompt}"
        )
        st.session_state.doc_chat_history.append(("assistant", answer))
        st.rerun()

